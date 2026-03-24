import os
import re
import base64
import hashlib
import tempfile
import subprocess
from typing import List, Dict, Any
import zipfile
from xml.etree import ElementTree

from pptx import Presentation

from .base import FileProcessor

from unstructured_inference.logger import logger
from unstructured_inference.models import tables
from unstructured.partition.auto import partition


tables_agent = tables.tables_agent
TABLE_TRANSFORMER_MODEL_PATH = ""

def custom_load_table_model():
    """Loads the Table agent."""

    if getattr(tables_agent, "model", None) is None:
        with tables_agent._lock:
            if getattr(tables_agent, "model", None) is None:
                logger.info("Loading the Table agent ...")
                print("path234: ", TABLE_TRANSFORMER_MODEL_PATH)
                tables_agent.initialize(TABLE_TRANSFORMER_MODEL_PATH)

    return

tables.load_agent = lambda: custom_load_table_model()


class UniversalImageExtractor(FileProcessor):
    """
    Multi-format image extractor for PDF, PPT, Excel, and Word.
    Uses LibreOffice for conversion when needed and reuses PDF extraction logic.
    """

    @staticmethod
    def _hash(data: bytes) -> str:
        return hashlib.md5(data).hexdigest()


    def _write_temp_file(self, data: bytes, suffix: str) -> str:
        tmp = tempfile.NamedTemporaryFile(delete=False, suffix=suffix)
        tmp.write(data)
        tmp.close()
        return tmp.name

    @staticmethod
    def detect_image_format(image_bytes: bytes) -> str:
        if image_bytes.startswith(b"\x89PNG"):
            return "png"
        elif image_bytes.startswith(b"\xFF\xD8\xFF"):
            return "jpg"
        else:
            return "png"


    def _convert_file(self, input_path: str, target_format: str) -> str:
    
        """
        Convert a file to the target format using LibreOffice.

        Args:
            input_path: Source file path.
            target_format: Target format, e.g. "pdf", "pptx", "xlsx".

        Returns:
            Output file path.
        """
        out_dir = os.path.dirname(input_path)

        cmd = [
            "soffice",
            "--headless",
            "--invisible",  # Ensure fully headless conversion.
            "--convert-to", f"{target_format}",
            input_path,
            "--outdir", out_dir
        ]

        try:
            subprocess.run(
                cmd,
                check=True,
                stdout=subprocess.DEVNULL,
                stderr=subprocess.DEVNULL,
                timeout=60  # Prevent hanging conversions.
            )

            base_name = os.path.splitext(input_path)[0]
            new_suffix = f".{target_format}"
            output_path = base_name + new_suffix

            if os.path.exists(output_path):
                return output_path
            else:
                raise FileNotFoundError(
                    f"Conversion failed: Output file {output_path} not found.")

        except subprocess.CalledProcessError as e:
            raise RuntimeError(
                f"LibreOffice conversion failed for {input_path}: {e}")
        except subprocess.TimeoutExpired:
            raise RuntimeError(
                f"LibreOffice conversion timed out for {input_path}")


    def _extract_pdf(self, pdf_path: str, **params) -> List[Dict]:
        table_transformer_model_path = params.get("table_transformer_model_path")
        unstructured_default_model_initialize_params_json_path = params.get(
            "unstructured_default_model_initialize_params_json_path"
        )
        if not table_transformer_model_path or not unstructured_default_model_initialize_params_json_path:
            return []
        global TABLE_TRANSFORMER_MODEL_PATH
        TABLE_TRANSFORMER_MODEL_PATH = table_transformer_model_path

        results = []
        seen = set()

        elements = partition(
            filename=pdf_path,
            strategy="hi_res",
            extract_images_in_pdf=True,
            extract_image_block_to_payload=True,
        )

        for el in elements:
            b64 = getattr(el.metadata, "image_base64", None)
            if not b64:
                continue

            img_bytes = base64.b64decode(b64)
            h = self._hash(img_bytes)
            if h in seen:
                continue
            seen.add(h)

            coords = getattr(el.metadata, "coordinates", None)
            coord_dict = None

            if coords and hasattr(coords, 'points') and coords.points:
                pts = coords.points  # tuple of (x,y)
                xs = [p[0] for p in pts]
                ys = [p[1] for p in pts]
                coord_dict = {
                    "x1": min(xs),
                    "y1": min(ys),
                    "x2": max(xs),
                    "y2": max(ys),
                }

            page_num = getattr(el.metadata, "page_number", None)

            results.append({
                "position": {
                    "page_number": page_num,
                    "coordinates": coord_dict
                },
                "image_format": self.detect_image_format(img_bytes),
                "image_bytes": img_bytes
            })

        return results


    def _extract_excel(self, xlsx_path):
        results = []
        seen = set()

        with zipfile.ZipFile(xlsx_path) as z:
            ns = {
                "xdr": "http://schemas.openxmlformats.org/drawingml/2006/spreadsheetDrawing",
                "a": "http://schemas.openxmlformats.org/drawingml/2006/main",
                "r": "http://schemas.openxmlformats.org/officeDocument/2006/relationships",
            }

            workbook = ElementTree.fromstring(z.read("xl/workbook.xml"))
            sheets = {}
            for s in workbook.findall(".//{http://schemas.openxmlformats.org/spreadsheetml/2006/main}sheet"):
                sheets[s.get("r:id")] = s.get("name")

            sheet_files = [f for f in z.namelist(
            ) if f.startswith("xl/worksheets/sheet")]

            for sheet_file in sheet_files:
                sheet_xml = ElementTree.fromstring(z.read(sheet_file))
                drawing = sheet_xml.find(
                    ".//{http://schemas.openxmlformats.org/spreadsheetml/2006/main}drawing")

                if drawing is None:
                    continue

                rId = drawing.get(
                    "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id")
                rel_path = sheet_file.replace(
                    "worksheets", "worksheets/_rels") + ".rels"

                if rel_path not in z.namelist():
                    continue

                rel_xml = ElementTree.fromstring(z.read(rel_path))
                drawing_file = None

                for r in rel_xml:
                    if r.get("Id") == rId:
                        drawing_file = "xl/" + \
                            r.get("Target").replace("../", "")
                        break

                if drawing_file is None:
                    continue

                sheet_name = os.path.basename(sheet_file)
                drawing_root = ElementTree.fromstring(z.read(drawing_file))

                rel_file = drawing_file.replace(
                    "drawings/", "drawings/_rels/") + ".rels"
                if rel_file not in z.namelist():
                    continue

                rel_root = ElementTree.fromstring(z.read(rel_file))
                rel_map = {
                    r.get("Id"): "xl/" + r.get("Target").replace("../", "")
                    for r in rel_root
                }

                anchors = drawing_root.findall(".//xdr:twoCellAnchor", ns) + \
                    drawing_root.findall(".//xdr:oneCellAnchor", ns)

                for anchor in anchors:
                    from_node = anchor.find("xdr:from", ns)
                    if from_node is None:
                        continue

                    row1 = int(from_node.find("xdr:row", ns).text) + 1
                    col1 = int(from_node.find("xdr:col", ns).text) + 1

                    to_node = anchor.find("xdr:to", ns)
                    if to_node is not None:
                        row2 = int(to_node.find("xdr:row", ns).text) + 1
                        col2 = int(to_node.find("xdr:col", ns).text) + 1
                    else:
                        row2, col2 = row1, col1

                    blip = anchor.find(".//a:blip", ns)
                    if blip is None:
                        continue

                    rId = blip.get(
                        "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}embed")
                    if rId not in rel_map:
                        continue

                    img_bytes = z.read(rel_map[rId])
                    h = self._hash(img_bytes)

                    if h in seen:
                        continue
                    seen.add(h)

                    results.append({
                        "position": {
                            "sheet_name": sheet_name,
                            "coordinates": {
                                "x1": col1,
                                "x2": col2,
                                "y1": row1,
                                "y2": row2
                            }
                        },
                        "image_format": self.detect_image_format(img_bytes),
                        "image_bytes": img_bytes
                    })

        return results


    def _extract_pptx(self, pptx_path: str, **params) -> List[Dict]:
        if Presentation is None:
            raise RuntimeError("python-pptx is required to extract images from PPTX files.")
        prs = Presentation(pptx_path)
        results = []
        seen = set()
        emu_per_inch = params.get("emu_per_inch", 914400)
        dpi = params.get("dpi", 96)
        
        def _emu_to_px(emu: int, emu_per_inch: int, dpi: int) -> int:
            return int((emu / emu_per_inch) * dpi)
        

        slide_w = _emu_to_px(prs.slide_width, emu_per_inch, dpi)
        slide_h = _emu_to_px(prs.slide_height, emu_per_inch, dpi)

        for slide_index, slide in enumerate(prs.slides):
            for shape in slide.shapes:
                if not hasattr(shape, "image"):
                    continue

                img_bytes = shape.image.blob
                h = self._hash(img_bytes)
                if h in seen:
                    continue
                seen.add(h)

                x = _emu_to_px(shape.left, emu_per_inch, dpi)
                y = _emu_to_px(shape.top, emu_per_inch, dpi)
                w = _emu_to_px(shape.width, emu_per_inch, dpi)
                h_px = _emu_to_px(shape.height, emu_per_inch, dpi)

                results.append({
                    "position": {
                        "page_number": slide_index + 1,
                        "coordinates": {
                            "x1": x,
                            "y1": y,
                            "x2": x + w,
                            "y2": y + h_px,
                            "slide_width": slide_w,
                            "slide_height": slide_h,
                        },
                    },
                    "image_format": self.detect_image_format(img_bytes),
                    "image_bytes": img_bytes
                })

        return results


    def process_file(self, file_bytes: bytes, chunking_strategy: str, filename: str, **params) -> List[Dict[str, Any]]:
        suffix = os.path.splitext(filename)[1].lower()
        temp_path = self._write_temp_file(file_bytes, suffix)
        converted_path = None

        try:
            if suffix == ".xlsx":
                return self._extract_excel(temp_path)
            if suffix == ".xls":
                converted_path = self._convert_file(temp_path, "xlsx")
                return self._extract_excel(converted_path)

            if suffix == ".pptx":
                return self._extract_pptx(temp_path, **params)
            if suffix == ".ppt":
                converted_path = self._convert_file(temp_path, "pptx")
                return self._extract_pptx(converted_path, **params)

            if suffix in [".docx", ".doc"]:
                converted_path = self._convert_file(temp_path, "pdf")
                return self._extract_pdf(converted_path, **params)

            if suffix == ".pdf":
                return self._extract_pdf(temp_path, **params)

            return []

        finally:
            files_to_clean = [temp_path]
            if converted_path and os.path.exists(converted_path):
                files_to_clean.append(converted_path)

            base = os.path.splitext(temp_path)[0]
            for ext in [".docx", ".pptx", ".xlsx", ".pdf"]:
                potential_file = base + ext
                if potential_file != converted_path and potential_file != temp_path:
                    files_to_clean.append(potential_file)

            for f_path in files_to_clean:
                if f_path and os.path.exists(f_path):
                    try:
                        os.remove(f_path)
                    except Exception:
                        pass


if __name__ == "__main__":
    extractor = UniversalImageExtractor()

    input_path = r"C:\Users\pc\Desktop\files\docx.docx"

    output_dir = "output_images"
    os.makedirs(output_dir, exist_ok=True)

    if not os.path.exists(input_path):
        print(f"Error: Input file not found: {input_path}")
    else:
        with open(input_path, "rb") as f:
            file_bytes = f.read()

        images = extractor.process_file(
            file_bytes, os.path.basename(input_path))
        if not images:
            print("No images found.")
        else:
            for i, img_info in enumerate(images, start=1):
                img_data = img_info["image_bytes"]
                img_fmt = img_info.get("image_format", "png")
                img_filename = f"image_{i}.{img_fmt}"
                img_path = os.path.join(output_dir, img_filename)

                with open(img_path, "wb") as f:
                    f.write(img_data)

                print(f"Saved: {img_path} (Format: {img_fmt})")

            print(f"\nFinished. Total {len(images)} images extracted.")
